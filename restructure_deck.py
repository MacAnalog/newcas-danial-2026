"""Restructure the NEWCAS deck into a 15-min, paper-focused talk build.

- Reuses Danial's existing on-topic slides (title, bio, challenges, score-shaping
  core, tool, future, conclusion, thanks).
- Inserts 3 NEW plot slides (convergence / trade-off / distributions) with
  takeaway titles + speaker notes.
- Reorders + subsets the slide list to the 14-slide arc from TALK_OUTLINE.md.

Original deck is never modified; output is a new file. Orphaned parts are removed
afterward by the pptx skill's clean.py + pack.py.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

SRC = "presentation/NEWCAS_ScoreShaping_with_SpiceXplorer.pptx"
OUT = "presentation/NEWCAS_ScoreShaping_TALK_15min.pptx"

INK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # 'Blank 2' has a dark burgundy background
BLANK_LAYOUT_NAME = "Blank 2"

# (existing 1-based slide number) in final order; "NEW:<key>" marks an inserted plot slide
FINAL_ORDER = [
    1,            # title
    4,            # speaker / bio
    7,            # problem: hard-to-quantify trade-offs
    19,           # framing: CSP + single-objective fitness
    20,           # score shaping: sigmoid normalization
    21,           # aggregation strategy / results context
    22,           # insight: linear stuck
    "NEW:conv",   # result 1 — convergence
    "NEW:trade",  # result 2 — trade-off outcome
    "NEW:dist",   # result 3 — exploration behaviour
    23,           # SpiceXplorer tool
    31,           # future: agentic
    34,           # conclusion
    36,           # thank you
]

PLOTS = {
    "conv": dict(
        img="plots/talk/convergence_trace.png",
        title="Same solver, same budget: sigmoid reaches feasibility, linear stalls",
        notes=("Result 1 of 3 — convergence. Both runs are Differential Evolution, 2000 sims, "
               "identical search space. Only the score SHAPE differs. Sigmoid crosses zero "
               "(all constraints met) at iter 1050 and then banks reward beyond spec. Linear "
               "never reaches feasibility — 0 fully-feasible points in 2000 sims. LHS and BFGS "
               "stall and are excluded for clarity. Punchline: the score function, not the "
               "solver, decided the outcome."),
    ),
    "trade": dict(
        img="plots/talk/tradeoff_clouds.png",
        title="Sigmoid meets every spec; linear blows the 25 \u00b5A current budget (43.7 \u00b5A)",
        notes=("Result 2 of 3 — the trade-off. Top-100 designs over the LHS cloud. "
               "Left/middle: current vs gain and vs UGF. Right: spec margin per metric. "
               "Sigmoid best = 20.6 \u00b5A, meets all specs. Linear best maximizes UGF (527 vs "
               "236 MHz is the WRONG lesson) and gain but lands at 43.7 \u00b5A, 1.7x over budget. "
               "This is the masking effect in numbers: the large MHz and dB magnitudes drowned "
               "the \u00b5A violation in the linear cost. Mirrors Table III."),
    ),
    "dist": dict(
        img="plots/talk/metric_distributions.png",
        title="Exploration: sigmoid keeps UGF diversity while pulling current below spec",
        notes=("Result 3 of 3 — exploration behaviour. Histograms of every visited design. "
               "Top row UGF, bottom row supply current; dashed line = spec target. Sigmoid "
               "spreads UGF broadly (no over-fit) and concentrates current below 25 \u00b5A. "
               "Linear over-concentrates at very high UGF and higher current. LHS stalls at the "
               "~5 \u00b5A sub-threshold reference bias. Sigmoid explores the USEFUL low-power "
               "region. Sets up the ~50%% power saving headline."),
    ),
}


def emu_fit(img_path, max_w_in, max_h_in):
    w, h = Image.open(img_path).size
    ar = w / h
    box_ar = max_w_in / max_h_in
    if ar >= box_ar:
        out_w = max_w_in
        out_h = max_w_in / ar
    else:
        out_h = max_h_in
        out_w = max_h_in * ar
    return out_w, out_h


def add_plot_slide(prs, blank_layout, img, title, notes):
    slide = prs.slides.add_slide(blank_layout)
    sw = prs.slide_width
    sh = prs.slide_height
    # title textbox
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35),
                                  sw - Inches(1.2), Inches(1.25))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = WHITE
    # image, top-aligned just below the title, horizontally centered
    area_top = 1.95
    max_w = (sw / 914400) - 1.2          # 0.6" margins
    max_h = (sh / 914400) - area_top - 0.5
    out_w, out_h = emu_fit(img, max_w, max_h)
    left = (sw - Inches(out_w)) / 2
    top = Inches(area_top)
    slide.shapes.add_picture(img, left, top, height=Inches(out_h))
    # speaker notes
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def main():
    prs = Presentation(SRC)
    blank = next(l for l in prs.slide_layouts if l.name == BLANK_LAYOUT_NAME)

    # map existing 1-based number -> its sldId element (current order is 1..N)
    sldIdLst = prs.slides._sldIdLst
    existing_ids = list(sldIdLst)  # in current order == slide 1..37
    num_to_sldId = {i + 1: existing_ids[i] for i in range(len(existing_ids))}

    # add the 3 new plot slides (appended); capture their sldId elements
    new_sldId = {}
    for key in ["conv", "trade", "dist"]:
        cfg = PLOTS[key]
        add_plot_slide(prs, blank, cfg["img"], cfg["title"], cfg["notes"])
        new_sldId[f"NEW:{key}"] = list(sldIdLst)[-1]

    # build the final ordered list of sldId elements
    ordered = []
    for item in FINAL_ORDER:
        if isinstance(item, int):
            ordered.append(num_to_sldId[item])
        else:
            ordered.append(new_sldId[item])

    # detach all sldId children, re-append in the desired order (drops the rest)
    for sid in list(sldIdLst):
        sldIdLst.remove(sid)
    for sid in ordered:
        sldIdLst.append(sid)

    prs.save(OUT)
    print("saved", OUT, "with", len(list(sldIdLst)), "slides")


if __name__ == "__main__":
    main()
